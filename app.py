import streamlit as st
from pptx import Presentation
from googletrans import Translator
from io import BytesIO
import time

st.set_page_config(page_title="Chinese to English PPT Translator", page_icon="🌏")

st.title("🌏 Chinese → English PowerPoint Translator")
st.write("Upload a `.pptx` file in Chinese, and this app will translate all text into English — for free!")

uploaded_file = st.file_uploader("📂 Upload your PowerPoint file", type=["pptx"])

if uploaded_file:
    try:
        prs = Presentation(uploaded_file)
        translator = Translator()
        total_shapes = sum(len(slide.shapes) for slide in prs.slides)
        progress = st.progress(0)
        count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        try:
                            translated = translator.translate(text, src='zh-cn', dest='en').text
                            shape.text = translated
                        except Exception as e:
                            shape.text = text + "\n(Translation failed)"
                    count += 1
                    progress.progress(count / total_shapes)
                    time.sleep(0.05)  # Prevents overloading Google Translate

        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        st.success("✅ Translation complete! Download your English PowerPoint below.")
        st.download_button(
            label="📥 Download Translated PPT",
            data=output,
            file_name="translated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        st.error(f"⚠️ Something went wrong: {e}")

else:
    st.info("Please upload a PPTX file to start.")
